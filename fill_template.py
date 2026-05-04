"""Fill SeniorProjectPPT_Template.pptx with senior-project content in place,
preserving the template's existing fonts, colors, shapes, and layouts.

All numbers come from experiments/results/final_report.md (90 runs).
"""
from copy import deepcopy
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

SRC = "SeniorProjectPPT_Template.pptx"
DST = "SeniorProjectPPT_Template.pptx"


# ---------- formatting helpers ----------

def _capture_run_format(tf):
    """Pull font info from the first non-empty run in a text frame."""
    for para in tf.paragraphs:
        for run in para.runs:
            if run.text or run.font.name or run.font.size:
                return {
                    "size": run.font.size,
                    "bold": run.font.bold,
                    "italic": run.font.italic,
                    "name": run.font.name,
                    "color_xml": _capture_color_xml(run),
                }
    return {}


def _capture_color_xml(run):
    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        return None
    fill = rPr.find(qn("a:solidFill"))
    if fill is None:
        return None
    return deepcopy(fill)


def _apply_run_format(run, fmt):
    if not fmt:
        return
    if fmt.get("size") is not None:
        run.font.size = fmt["size"]
    if fmt.get("bold") is not None:
        run.font.bold = fmt["bold"]
    if fmt.get("italic") is not None:
        run.font.italic = fmt["italic"]
    if fmt.get("name"):
        run.font.name = fmt["name"]
    if fmt.get("color_xml") is not None:
        rPr = run._r.get_or_add_rPr()
        existing = rPr.find(qn("a:solidFill"))
        if existing is not None:
            rPr.remove(existing)
        rPr.append(deepcopy(fmt["color_xml"]))


def set_text(shape, lines, *, size=None, bold=None, align=None):
    """Replace the text frame contents of `shape` with `lines`,
    preserving formatting captured from the existing first run.
    """
    tf = shape.text_frame
    fmt = _capture_run_format(tf)
    if size is not None:
        fmt["size"] = Pt(size)
    if bold is not None:
        fmt["bold"] = bold

    captured_align = tf.paragraphs[0].alignment if tf.paragraphs else None
    final_align = align if align is not None else captured_align

    if isinstance(lines, str):
        lines = [lines]

    txBody = tf._txBody
    for p in list(txBody.findall(qn("a:p"))):
        txBody.remove(p)

    for line in lines:
        p = tf.add_paragraph()
        if final_align is not None:
            p.alignment = final_align
        run = p.add_run()
        run.text = line
        _apply_run_format(run, fmt)


def set_paragraphs(shape, paragraphs, *, align=None):
    """Each entry is (text, overrides) — gives one paragraph with one run.
    overrides may include size, bold, italic, color_rgb.
    """
    tf = shape.text_frame
    base_fmt = _capture_run_format(tf)
    captured_align = tf.paragraphs[0].alignment if tf.paragraphs else None
    final_align = align if align is not None else captured_align

    txBody = tf._txBody
    for p in list(txBody.findall(qn("a:p"))):
        txBody.remove(p)

    for text, overrides in paragraphs:
        p = tf.add_paragraph()
        if final_align is not None:
            p.alignment = final_align
        run = p.add_run()
        run.text = text
        fmt = dict(base_fmt)
        if overrides:
            if "color_rgb" in overrides:
                fmt["color_xml"] = None
            fmt.update({k: v for k, v in overrides.items() if k != "color_rgb"})
        _apply_run_format(run, fmt)
        if overrides and "color_rgb" in overrides:
            run.font.color.rgb = overrides["color_rgb"]


def set_runs(shape, segments, *, align=None):
    """Replace the text frame with multiple runs in a single paragraph."""
    tf = shape.text_frame
    base_fmt = _capture_run_format(tf)
    captured_align = tf.paragraphs[0].alignment if tf.paragraphs else None
    final_align = align if align is not None else captured_align

    txBody = tf._txBody
    for p in list(txBody.findall(qn("a:p"))):
        txBody.remove(p)

    p = tf.add_paragraph()
    if final_align is not None:
        p.alignment = final_align

    for text, overrides in segments:
        run = p.add_run()
        run.text = text
        fmt = dict(base_fmt)
        if overrides:
            if "color_rgb" in overrides:
                fmt["color_xml"] = None
            fmt.update({k: v for k, v in overrides.items() if k != "color_rgb"})
        _apply_run_format(run, fmt)
        if overrides and "color_rgb" in overrides:
            run.font.color.rgb = overrides["color_rgb"]


# ---------- main fill logic ----------

prs = Presentation(SRC)
ACCENT = RGBColor(0x1A, 0x42, 0x8A)  # template's primary blue

# ============ SLIDE 1: Title ============
s = prs.slides[0]
title_tf = s.shapes[0].text_frame
title_fmt = _capture_run_format(title_tf)
captured_align = title_tf.paragraphs[0].alignment if title_tf.paragraphs else None
txBody = title_tf._txBody
for p in list(txBody.findall(qn("a:p"))):
    txBody.remove(p)

p1 = title_tf.add_paragraph()
if captured_align is not None:
    p1.alignment = captured_align
r1 = p1.add_run()
r1.text = "How Good Are LLMs as Data Analysts?"
_apply_run_format(r1, title_fmt)

p2 = title_tf.add_paragraph()
if captured_align is not None:
    p2.alignment = captured_align
r2 = p2.add_run()
r2.text = "A Comparative Evaluation of Agentic Architectures on Real-World Survey Data"
sub_fmt = dict(title_fmt)
sub_fmt["size"] = Pt(20)
sub_fmt["bold"] = False
_apply_run_format(r2, sub_fmt)

# Names
set_text(s.shapes[1], "Ali Ibrahim   ·   Elsayed Azab   ·   Mohammed Sharaf", size=24)
# Supervisor + course
set_text(
    s.shapes[2],
    "Supervisor: Dr. Omar Alomeir   ·   CS 499 Senior Project   ·   Prince Sultan University   ·   Spring 2026",
    size=18,
)

# ============ SLIDE 2: Three Agent Architectures ============
s = prs.slides[1]
set_text(s.shapes[1], "Three Agent Architectures")
set_text(
    s.shapes[0],
    "Same LLM. Same toolset. Three designs — and a 37-point accuracy gap.",
)

# Box order on slide L→R: shape[6] Single, shape[5] Middle, shape[4] Right.
# Re-order story: Single → RAG → Multi-Agent.
def fill_agent_box(shape, name, description, metric):
    set_paragraphs(shape, [
        (name, {"size": Pt(18), "bold": True, "color_rgb": ACCENT}),
        (" ", {"size": Pt(6)}),
        (description, {"size": Pt(12), "bold": False, "color_rgb": ACCENT}),
        (" ", {"size": Pt(10)}),
        (metric, {"size": Pt(22), "bold": True, "color_rgb": ACCENT}),
    ], align=PP_ALIGN.CENTER)

fill_agent_box(
    s.shapes[6],
    "Single Agent",
    "ReAct loop. One LLM iterates with shared tools.",
    "42%  ·  53s",
)
fill_agent_box(
    s.shapes[5],
    "RAG Agent",
    "Retrieves codebook context to ground the prompt.",
    "62%  ·  54s",
)
fill_agent_box(
    s.shapes[4],
    "Multi-Agent",
    "Planner → Analyst → Reviewer with retry loop.",
    "79%  ·  166s",
)

# ============ SLIDE 3: Contributions ============
s = prs.slides[2]
set_text(s.shapes[1], "Contributions", size=32, bold=True)
set_text(s.shapes[0], "Four contributions of this senior project")

contributions = [
    (s.shapes[2],
     "Reproducible end-to-end pipeline: 30 questions × 3 architectures × 3 datasets = 90 runs"),
    (s.shapes[3],
     "Direct head-to-head comparison of Single, Multi-Agent, and RAG architectures on real survey data"),
    (s.shapes[4],
     "Empirical finding: architecture lifts accuracy by +37 points on the same LLM and toolset"),
    (s.shapes[5],
     "Open-source agents, evaluation framework, ground-truth library, and Flask comparison dashboard"),
]
for shape, text in contributions:
    set_paragraphs(shape, [
        (text, {"size": Pt(14), "bold": True, "color_rgb": ACCENT}),
    ], align=PP_ALIGN.LEFT)

# ============ SLIDE 4: System Architecture ============
s = prs.slides[3]
set_text(s.shapes[1], "System Architecture")
# Place the existing TextBox 11 at the top of the dark rectangle as the section header
tb = s.shapes[5]
tb.left = Inches(4.7)
tb.top = Inches(1.95)
tb.width = Inches(7.6)
tb.height = Inches(0.6)
WHITE_RGB = RGBColor(0xFF, 0xFF, 0xFF)
set_paragraphs(tb, [
    ("Modular pipeline — three agents over a shared toolset.",
        {"size": Pt(16), "bold": True, "color_rgb": WHITE_RGB}),
])

# ============ SLIDE 5: Evaluation Pipeline ============
s = prs.slides[4]
set_text(s.shapes[1], "Evaluation Pipeline")
set_text(s.shapes[0], "Seven stages — from question to comparison report")
# Clear the giant overlapping placeholder textbox
set_text(s.shapes[29], "")

steps = [
    (2,  "Curate 30 analytical questions across GSS, Arab Barometer, and WVS"),
    (6,  "Implement expert ground-truth functions in Python with survey weights"),
    (10, "Build Single, Multi-Agent, and RAG architectures over a shared tool layer"),
    (14, "Index codebook PDFs into ChromaDB with all-MiniLM-L6-v2 embeddings"),
    (30, "Run every agent on every question with a 120-second sandbox timeout"),
    (21, "Score on accuracy, completeness, weight usage, latency, and retries"),
    (25, "Aggregate per-agent results and generate the markdown comparison report"),
]
for shape_idx, text in steps:
    sh = s.shapes[shape_idx]
    tf = sh.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.1)
    txBody = tf._txBody
    for p in list(txBody.findall(qn("a:p"))):
        txBody.remove(p)
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.name = "Somar Sans"
    run.font.color.rgb = ACCENT

# ============ SLIDE 6: Key Findings ============
s = prs.slides[5]
set_text(s.shapes[1], "Key Findings")
set_paragraphs(s.shapes[2], [
    ("Architecture > model. Same LLM, +37 points of accuracy: Multi-Agent reaches 79% vs Single at 42%.",
        {"size": Pt(13.5), "bold": True}),
    ("", {"size": Pt(4)}),
    ("Verification is the largest single lever — the Reviewer step lifts weight usage from 60% to 100%.",
        {"size": Pt(13)}),
    ("", {"size": Pt(4)}),
    ("RAG retrieval gives +20 points at near-zero latency cost (54s vs 53s for Single).",
        {"size": Pt(13)}),
    ("", {"size": Pt(4)}),
    ("Multi-Agent quality has a price: ~3× latency (166s) and 0.83 review retries on average.",
        {"size": Pt(13)}),
    ("", {"size": Pt(4)}),
    ("All architectures excel on Arab Barometer; WVS remains the hardest dataset for every agent.",
        {"size": Pt(13)}),
])

# ============ SLIDE 7: Performance by Architecture ============
s = prs.slides[6]
set_text(s.shapes[0], "Performance by Architecture")

# Re-order labels to tell the improvement story top→bottom: Single → RAG → Multi.
group_labels = [
    (s.shapes[2], "Single Agent"),
    (s.shapes[3], "RAG Agent"),
    (s.shapes[4], "Multi-Agent"),
]
for grp, label in group_labels:
    label_shape = grp.shapes[0]
    set_text(label_shape, label)

# Right-side metric bars — REAL numbers from final_report.md (90 runs).
# Single → 42% acc, 53s avg, 60% weight, 0% errors
set_runs(s.shapes[1], [
    ("42%", {"bold": True, "color_rgb": ACCENT}),
    (" accuracy   ·   ", None),
    ("53s", {"bold": True}),
    (" avg time   ·   ", None),
    ("60%", {"bold": True}),
    (" weight usage   ·   ", None),
    ("0%", {"bold": True}),
    (" errors", None),
])
# RAG → 62% acc, 54s avg, 70% weight, 0% errors
set_runs(s.shapes[7], [
    ("62%", {"bold": True, "color_rgb": ACCENT}),
    (" accuracy   ·   ", None),
    ("54s", {"bold": True}),
    (" avg time   ·   ", None),
    ("70%", {"bold": True}),
    (" weight usage   ·   ", None),
    ("0%", {"bold": True}),
    (" errors", None),
])
# Multi → 79% acc, 166s avg, 100% weight, 0% errors
set_runs(s.shapes[8], [
    ("79%", {"bold": True, "color_rgb": ACCENT}),
    (" accuracy   ·   ", None),
    ("166s", {"bold": True}),
    (" avg time   ·   ", None),
    ("100%", {"bold": True}),
    (" weight usage   ·   ", None),
    ("0.83", {"bold": True}),
    (" avg retries", None),
])

# ============ SLIDE 8: Thank you ============
s = prs.slides[7]
# Expand title placeholder so multiple lines fit comfortably and breathe
sh = s.shapes[0]
sh.left = Inches(0.6)
sh.top = Inches(2.0)
sh.width = Inches(12.13)
sh.height = Inches(3.6)
set_paragraphs(sh, [
    ("Thank You", {"size": Pt(64), "bold": True}),
    (" ", {"size": Pt(14)}),
    ("Questions & Discussion", {"size": Pt(24), "bold": False}),
    (" ", {"size": Pt(18)}),
    ("Ali Ibrahim   ·   Elsayed Azab   ·   Mohammed Sharaf",
        {"size": Pt(16), "bold": False}),
    (" ", {"size": Pt(6)}),
    ("Supervisor: Dr. Omar Alomeir   ·   Prince Sultan University   ·   Spring 2026",
        {"size": Pt(13), "bold": False}),
], align=PP_ALIGN.CENTER)

# ============================================================
# VISUAL ENHANCEMENTS — overlay shapes for richer look
# (idempotent: every shape is named ENH_* and removed before re-adding)
# ============================================================

CYAN = RGBColor(0x29, 0xAB, 0xE2)        # template's bright cyan
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCD, 0xDB, 0xEC)       # light blue tint
TRACK = RGBColor(0xE8, 0xEE, 0xF6)       # progress-bar track


def remove_enhancements(slide):
    spTree = slide.shapes._spTree
    for shape in list(slide.shapes):
        if shape.name and shape.name.startswith("ENH_"):
            spTree.remove(shape._element)


def _set_name(shape, name):
    cnv = shape._element.find(".//" + qn("p:cNvPr"))
    if cnv is not None:
        cnv.set("name", name)


def add_rect(slide, x, y, w, h, *, fill=None, name="ENH_rect"):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is not None:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    _set_name(sh, name)
    return sh


def add_txt(slide, x, y, w, h, content, *, size=14, bold=False, italic=False,
            color=None, align=PP_ALIGN.LEFT, font=None,
            anchor=MSO_ANCHOR.TOP, name="ENH_txt"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = content
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    if color is not None:
        r.font.color.rgb = color
    if font:
        r.font.name = font
    _set_name(tb, name)
    return tb


def add_bar(slide, x, y, w, h, value, *, max_val=1.0,
            fill=None, track=None, name="ENH_bar"):
    fill = fill or ACCENT
    track = track or TRACK
    add_rect(slide, x, y, w, h, fill=track, name=name + "_t")
    fw_emu = int(w * (value / max_val))
    if fw_emu > 0:
        add_rect(slide, x, y, fw_emu, h, fill=fill, name=name + "_f")


# ----- Slide 1: subtle keyword strip beneath the title --------------------
s = prs.slides[0]
remove_enhancements(s)
# Faint cyan accent rule under the names line
add_rect(s, Inches(0.69), Inches(4.45), Inches(0.6), Emu(25000),
         fill=CYAN, name="ENH_s1_rule")
add_txt(s, Inches(0.69), Inches(4.55), Inches(11.5), Inches(0.35),
        "GSS  ·  ARAB BAROMETER VIII  ·  WORLD VALUES SURVEY 7",
        size=11, bold=True, color=WHITE, name="ENH_s1_keywords")
add_txt(s, Inches(0.69), Inches(4.85), Inches(11.5), Inches(0.35),
        "30 questions   ·   3 architectures   ·   3 datasets   =   90 controlled runs",
        size=11, italic=True, color=LIGHT, name="ENH_s1_meta")


# ----- Slide 3: big "01-04" badges before each contribution row -----------
s = prs.slides[2]
remove_enhancements(s)
# Row Y midpoints (approximate centers of each freeform)
row_centers = [3.63, 4.50, 5.45, 6.46]
for i, ymid in enumerate(row_centers, start=1):
    digit = f"{i:02d}"
    # Big translucent number on the right edge of each row
    add_txt(s, Inches(11.3), Inches(ymid - 0.45), Inches(0.95), Inches(0.9),
            digit, size=44, bold=True, color=CYAN,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
            name=f"ENH_s3_num_{i}")


# ----- Slide 4: architecture flow pills inside dark rectangle -------------
s = prs.slides[3]
remove_enhancements(s)
pills = [
    ("01  INPUT",      "Natural-language question + dataset selector"),
    ("02  AGENTS",     "Single (ReAct)   ·   Multi-Agent   ·   RAG"),
    ("03  TOOLS",      "load_dataset  ·  get_schema  ·  get_variable_info  ·  run_analysis_code"),
    ("04  DATA",       "GSS   ·   Arab Barometer Wave VIII   ·   World Values Survey 7"),
    ("05  EVALUATION", "ground-truth functions  ·  scoring metrics  ·  comparator"),
]
pill_x = 4.85
pill_w = 7.30
y0 = 2.70
gap = 0.10
pill_h = 0.58
for i, (tag, desc) in enumerate(pills):
    py = y0 + i * (pill_h + gap)
    add_txt(s, Inches(pill_x), Inches(py), Inches(2.0), Inches(pill_h),
            tag, size=11, bold=True, color=CYAN,
            anchor=MSO_ANCHOR.MIDDLE, name=f"ENH_s4_tag_{i}")
    add_txt(s, Inches(pill_x + 2.0), Inches(py), Inches(pill_w - 2.0), Inches(pill_h),
            desc, size=12, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE, name=f"ENH_s4_desc_{i}")
    if i < len(pills) - 1:
        add_rect(s, Inches(pill_x), Inches(py + pill_h + gap/2 - 0.01),
                 Inches(pill_w), Emu(7000), fill=CYAN,
                 name=f"ENH_s4_rule_{i}")


# ----- Slide 6: replace dense bullets with 4 big stat tiles ---------------
s = prs.slides[5]
remove_enhancements(s)
# Clear the rectangle's text — we overlay tiles instead
set_text(s.shapes[2], "")

# Stat tiles inside the dark-blue rectangle (L=4.59 T=2.79 W=7.51 H=2.85)
rect_x, rect_y, rect_w, rect_h = 4.59, 2.79, 7.51, 2.85
gap = 0.10
inner_pad = 0.18
n = 4
tile_w = (rect_w - 2 * inner_pad - (n - 1) * gap) / n
inner_x0 = rect_x + inner_pad

stats = [
    ("+37",  "PTS",  "Multi vs Single accuracy"),
    ("+20",  "PTS",  "RAG vs Single accuracy"),
    ("100%", "",     "Multi-Agent weight usage"),
    ("3.1×", "",     "Multi-Agent latency cost"),
]
for i, (number, unit, label) in enumerate(stats):
    tx = inner_x0 + i * (tile_w + gap)
    # Big number — anchored middle of upper region
    add_txt(s, Inches(tx), Inches(rect_y + 0.20), Inches(tile_w), Inches(1.05),
            number, size=40, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            name=f"ENH_s6_num_{i}")
    if unit:
        add_txt(s, Inches(tx), Inches(rect_y + 1.22), Inches(tile_w), Inches(0.25),
                unit, size=10, bold=True, color=CYAN,
                align=PP_ALIGN.CENTER, name=f"ENH_s6_unit_{i}")
    add_txt(s, Inches(tx), Inches(rect_y + 1.55), Inches(tile_w), Inches(0.6),
            label, size=10, color=LIGHT,
            align=PP_ALIGN.CENTER, name=f"ENH_s6_lbl_{i}")
    # tiny accent rule beneath the number
    add_rect(s, Inches(tx + tile_w / 2 - 0.18), Inches(rect_y + 1.42),
             Inches(0.36), Emu(20000), fill=CYAN, name=f"ENH_s6_rule_{i}")

# Bottom takeaway line
add_txt(s, Inches(rect_x + 0.2), Inches(rect_y + 2.30),
        Inches(rect_w - 0.4), Inches(0.45),
        "Same LLM. Same toolset.  ·  Architecture is the lever.",
        size=12, italic=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        name="ENH_s6_takeaway")


# ----- Slide 7: replace metric text with big number + progress bar --------
s = prs.slides[6]
remove_enhancements(s)

# Clear existing metric text — we render from scratch
for idx in (1, 7, 8):
    set_text(s.shapes[idx], "")

# Each row: big % + filled progress bar + small metrics text
ROWS = [
    (s.shapes[1], 42, "53s avg time   ·   60% weight   ·   0% errors"),
    (s.shapes[7], 62, "54s avg time   ·   70% weight   ·   0% errors"),
    (s.shapes[8], 79, "166s avg time   ·   100% weight   ·   0.83 retries"),
]
for shape, accuracy, metrics in ROWS:
    rx = shape.left / 914400
    ry = shape.top / 914400
    rh = shape.height / 914400

    # Big accuracy percent (left)
    add_txt(s, Inches(rx + 0.15), Inches(ry), Inches(1.4), Inches(rh),
            f"{accuracy}%",
            size=28, bold=True, color=ACCENT,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
            name=f"ENH_s7_acc_{accuracy}")

    # Progress bar (middle) — vertically centered
    bar_x = rx + 1.55
    bar_w = 3.6
    bar_h = 0.22
    bar_y = ry + (rh - bar_h) / 2
    add_bar(s, Inches(bar_x), Inches(bar_y), Inches(bar_w), Inches(bar_h),
            accuracy / 100.0, fill=CYAN, track=TRACK,
            name=f"ENH_s7_bar_{accuracy}")

    # Small metrics text (right)
    add_txt(s, Inches(bar_x + bar_w + 0.20), Inches(ry),
            Inches(9.21 - 1.55 - bar_w - 0.30), Inches(rh),
            metrics,
            size=11, color=ACCENT, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
            name=f"ENH_s7_metrics_{accuracy}")


prs.save(DST)
print(f"Saved: {DST}")
