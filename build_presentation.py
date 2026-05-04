"""Senior-project expo deck — dark, minimal, modern."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- Palette (dark, restrained, one accent) ----------
INK        = RGBColor(0x0B, 0x0E, 0x14)   # near-black background
SURFACE    = RGBColor(0x12, 0x16, 0x1F)   # slightly elevated
BORDER     = RGBColor(0x1F, 0x25, 0x32)   # divider line
TEXT       = RGBColor(0xEA, 0xEC, 0xEF)   # primary text
DIM        = RGBColor(0x9A, 0xA0, 0xA8)   # secondary text
MUTED      = RGBColor(0x5F, 0x66, 0x70)   # tertiary
ACCENT     = RGBColor(0x6E, 0xE7, 0xB7)   # mint
ACCENT_DIM = RGBColor(0x35, 0x82, 0x6B)   # darker mint
WARM       = RGBColor(0xFB, 0xBF, 0x24)   # warm contrast (used sparingly)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

MONO = "Consolas"
SANS = "Calibri"


# ---------------- Helpers ----------------
def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def alpha_fill(shape, color, alpha_pct):
    """alpha in percent (0-100)."""
    fill_solid(shape, color)
    srgb = shape.fill._xPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    etree.SubElement(srgb, qn('a:alpha'), attrib={'val': str(int(alpha_pct * 1000))})


def rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        if line_w:
            sh.line.width = line_w
    if fill is not None:
        fill_solid(sh, fill)
    return sh


def vline(slide, x, y, h, color=BORDER, width=Pt(0.75)):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x, y + h)
    line.line.color.rgb = color
    line.line.width = width
    return line


def hline(slide, x, y, w, color=BORDER, width=Pt(0.75)):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x + w, y)
    line.line.color.rgb = color
    line.line.width = width
    return line


def text(slide, x, y, w, h, content, *, size=14, bold=False, italic=False,
         color=TEXT, align=PP_ALIGN.LEFT, font=SANS,
         anchor=MSO_ANCHOR.TOP, tracking=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = content
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    if tracking:
        rPr = r._r.get_or_add_rPr()
        rPr.set('spc', str(tracking))
    return tb


def lines(slide, x, y, w, h, items, *, size=12, color=DIM, line_space=Pt(8), font=SANS):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = line_space
        if item == "":
            r = p.add_run()
            r.text = " "
            r.font.size = Pt(4)
            continue
        r = p.add_run()
        r.text = item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = font


# ---------------- Page chrome ----------------
def page_bg(slide):
    rect(slide, 0, 0, SW, SH, fill=INK)


def slide_meta(slide, num, kicker):
    # top-left mono kicker
    text(slide, Inches(0.6), Inches(0.55), Inches(8), Inches(0.3),
         f"{num:02d}  ·  {kicker.upper()}",
         size=10, color=ACCENT, font=MONO, tracking=200)
    # top-right brand
    text(slide, Inches(5.0), Inches(0.55), Inches(7.7), Inches(0.3),
         "LLM DATA SCIENTIST",
         size=10, color=MUTED, font=MONO, align=PP_ALIGN.RIGHT, tracking=300)
    # thin divider top
    hline(slide, Inches(0.6), Inches(0.95), Inches(12.13), color=BORDER, width=Pt(0.5))
    # thin divider bottom
    hline(slide, Inches(0.6), Inches(7.0), Inches(12.13), color=BORDER, width=Pt(0.5))
    # bottom-right page count
    text(slide, Inches(11.0), Inches(7.1), Inches(1.7), Inches(0.3),
         f"{num:02d} / 10",
         size=9, color=MUTED, font=MONO, align=PP_ALIGN.RIGHT)
    text(slide, Inches(0.6), Inches(7.1), Inches(8), Inches(0.3),
         "SENIOR PROJECT  ·  EXPO 2026",
         size=9, color=MUTED, font=MONO, tracking=200)


def slide_title(slide, title, subtitle=None):
    text(slide, Inches(0.6), Inches(1.3), Inches(12.13), Inches(1.0),
         title, size=40, bold=True, color=TEXT, font=SANS)
    if subtitle:
        text(slide, Inches(0.6), Inches(2.1), Inches(12.13), Inches(0.5),
             subtitle, size=15, color=DIM, font=SANS)


# ---------------- Custom bar chart (dark, flat) ----------------
def hbar_chart(slide, x, y, w, h, *, labels, values, max_value=None,
               accent=ACCENT, label_w=Inches(2.0), value_suffix="",
               value_fmt="{:.0f}"):
    """Horizontal bars with labels left and values right."""
    n = len(labels)
    row_h = h / n
    bar_area_w = w - label_w - Inches(1.0)
    if max_value is None:
        max_value = max(values) * 1.1
    track_h = Inches(0.16)
    for i, (lbl, val) in enumerate(zip(labels, values)):
        cy = y + row_h * i
        center_y = cy + row_h / 2
        # label
        text(slide, x, cy, label_w, row_h, lbl,
             size=12, bold=True, color=TEXT, font=SANS,
             anchor=MSO_ANCHOR.MIDDLE)
        # track (background bar)
        track = rect(slide, x + label_w, center_y - track_h / 2,
                     bar_area_w, track_h, fill=SURFACE)
        # value bar
        bar_w = int(bar_area_w * (val / max_value))
        if bar_w > 0:
            rect(slide, x + label_w, center_y - track_h / 2,
                 bar_w, track_h, fill=accent)
        # value label at right
        text(slide, x + label_w + bar_area_w + Inches(0.1), cy,
             Inches(0.9), row_h,
             value_fmt.format(val) + value_suffix,
             size=12, bold=True, color=accent, font=MONO,
             anchor=MSO_ANCHOR.MIDDLE)


def grouped_hbar(slide, x, y, w, h, *, categories, series):
    """Three series compared per category. series: list of (name, color, values)."""
    n_cat = len(categories)
    n_ser = len(series)
    cat_h = h / n_cat
    bar_h = Inches(0.18)
    gap = Inches(0.06)
    label_w = Inches(2.0)
    bar_area_w = w - label_w - Inches(0.7)
    max_val = 1.0  # accuracy 0-1
    for i, cat in enumerate(categories):
        cy = y + cat_h * i + Inches(0.1)
        # category label
        text(slide, x, cy, label_w, cat_h - Inches(0.2), cat,
             size=12, bold=True, color=TEXT, font=SANS, anchor=MSO_ANCHOR.MIDDLE)
        # bars stacked vertically per series within category
        ser_total_h = n_ser * bar_h + (n_ser - 1) * gap
        ser_y0 = cy + (cat_h - Inches(0.2) - ser_total_h) / 2
        for j, (sname, color, vals) in enumerate(series):
            v = vals[i]
            by = ser_y0 + j * (bar_h + gap)
            # track
            rect(slide, x + label_w, by, bar_area_w, bar_h, fill=SURFACE)
            bw = int(bar_area_w * (v / max_val))
            if bw > 0:
                rect(slide, x + label_w, by, bw, bar_h, fill=color)
            # value
            text(slide, x + label_w + bar_area_w + Inches(0.05), by - Inches(0.02),
                 Inches(0.65), bar_h + Inches(0.04),
                 f"{int(v*100)}%",
                 size=10, bold=True, color=color, font=MONO,
                 anchor=MSO_ANCHOR.MIDDLE)


# =============================================================
# SLIDE 1 — Title
# =============================================================
s = prs.slides.add_slide(BLANK)
page_bg(s)

# Decorative giant watermark "01"
text(s, Inches(9.5), Inches(0.6), Inches(3.5), Inches(2.0),
     "01", size=130, bold=True, color=SURFACE, font=MONO, align=PP_ALIGN.RIGHT)

# Top kicker
text(s, Inches(0.7), Inches(0.7), Inches(8), Inches(0.3),
     "SENIOR PROJECT  ·  EXPO 2026",
     size=11, color=ACCENT, font=MONO, tracking=300)

# Tiny accent dot
dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(2.85), Inches(0.12), Inches(0.12))
fill_solid(dot, ACCENT); dot.line.fill.background()

# Hero title — generous, big
text(s, Inches(0.7), Inches(3.05), Inches(12), Inches(1.5),
     "LLM Data Scientist.",
     size=84, bold=True, color=TEXT, font=SANS)

# Tagline
text(s, Inches(0.7), Inches(4.55), Inches(12), Inches(0.6),
     "Evaluating large language models as autonomous data analysts.",
     size=22, italic=True, color=DIM)

# Thin accent rule
rect(s, Inches(0.7), Inches(5.6), Inches(0.5), Emu(20000), fill=ACCENT)

# Author + dataset line
text(s, Inches(0.7), Inches(5.75), Inches(8), Inches(0.4),
     "Sayed Baraka", size=16, bold=True, color=TEXT)
text(s, Inches(0.7), Inches(6.1), Inches(11), Inches(0.4),
     "Computer Science  ·  2026", size=11, color=DIM, font=MONO)
text(s, Inches(0.7), Inches(6.4), Inches(11), Inches(0.4),
     "GSS  ·  WORLD VALUES SURVEY 7  ·  ARAB BAROMETER VIII",
     size=10, color=MUTED, font=MONO, tracking=300)

# bottom rule + page
hline(s, Inches(0.6), Inches(7.0), Inches(12.13), color=BORDER, width=Pt(0.5))
text(s, Inches(11.0), Inches(7.1), Inches(1.7), Inches(0.3),
     "01 / 10", size=9, color=MUTED, font=MONO, align=PP_ALIGN.RIGHT)


# =============================================================
# SLIDE 2 — Problem
# =============================================================
s = prs.slides.add_slide(BLANK)
page_bg(s)
slide_meta(s, 2, "Problem")

# Hero question
text(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.0),
     "Can LLMs really analyze data?",
     size=44, bold=True, color=TEXT)

text(s, Inches(0.6), Inches(2.6), Inches(11.5), Inches(0.6),
     "Survey analysis is weighted, codebook-driven, and methodologically strict —",
     size=16, color=DIM)
text(s, Inches(0.6), Inches(2.95), Inches(11.5), Inches(0.6),
     "far beyond the toy datasets in existing LLM benchmarks.",
     size=16, color=DIM)

# Four research questions in a clean horizontal row
hline(s, Inches(0.6), Inches(4.6), Inches(12.13), color=BORDER, width=Pt(0.5))
text(s, Inches(0.6), Inches(4.25), Inches(8), Inches(0.3),
     "RESEARCH QUESTIONS",
     size=10, color=ACCENT, font=MONO, tracking=300)

rqs = [
    ("RQ1", "Accuracy",
     "How accurately can LLM agents answer\nanalytical questions on survey data?"),
    ("RQ2", "Architecture",
     "Does multi-agent decomposition outperform\na single ReAct agent?"),
    ("RQ3", "Grounding",
     "Does codebook retrieval (RAG) improve\ncorrectness on coded variables?"),
    ("RQ4", "Trade-offs",
     "What is the accuracy / latency / cost\ntrade-off across architectures?"),
]
col_w = Inches(2.95)
gap = Inches(0.13)
start_x = Inches(0.6)
for i, (tag, head, body) in enumerate(rqs):
    cx = start_x + i * (col_w + gap)
    text(s, cx, Inches(4.85), col_w, Inches(0.3),
         tag, size=10, color=ACCENT, font=MONO, tracking=200)
    text(s, cx, Inches(5.15), col_w, Inches(0.5),
         head, size=20, bold=True, color=TEXT)
    text(s, cx, Inches(5.7), col_w, Inches(1.2),
         body, size=11.5, color=DIM)
    # vertical separator (except last)
    if i < 3:
        vline(s, cx + col_w + gap / 2, Inches(4.85), Inches(2.0), color=BORDER, width=Pt(0.5))


# =============================================================
# SLIDE 3 — Datasets
# =============================================================
s = prs.slides.add_slide(BLANK)
page_bg(s)
slide_meta(s, 3, "Datasets")

slide_title(s, "Three real-world surveys.",
            "~190,000 respondents · 60+ countries · weighted, coded, multi-format.")

# Big stat row
stats = [
    ("190K+", "RESPONDENTS"),
    ("60+", "COUNTRIES"),
    ("3", "FORMATS"),
    ("100%", "WEIGHTED"),
]
sx0 = Inches(0.6)
sw = Inches(2.95)
sgap = Inches(0.13)
for i, (val, lbl) in enumerate(stats):
    x = sx0 + i * (sw + sgap)
    text(s, x, Inches(2.85), sw, Inches(0.9),
         val, size=44, bold=True, color=ACCENT, font=MONO)
    text(s, x, Inches(3.75), sw, Inches(0.3),
         lbl, size=10, color=DIM, font=MONO, tracking=300)

# Divider
hline(s, Inches(0.6), Inches(4.45), Inches(12.13), color=BORDER, width=Pt(0.5))

# Three dataset blocks (no boxes, just typography)
ds = [
    ("GSS", "General Social Survey", "1972 – 2022",
     ["~72,000 respondents", "Stata · 567 MB", "Weight  ›  WTSSPS",
      "U.S. attitudes, demographics, behaviors."]),
    ("WVS 7", "World Values Survey", "Wave 7  ·  2017 – 2022",
     ["~95,000 respondents · 64 countries", "CSV · 182 MB", "Weight  ›  W_WEIGHT",
      "Cross-national values & beliefs."]),
    ("AB VIII", "Arab Barometer", "Wave VIII  ·  2023 – 2024",
     ["~26,000 respondents · MENA region", "CSV / DTA / SAV", "Weight  ›  WT",
      "Public opinion & governance."]),
]
col_w = Inches(3.95)
col_gap = Inches(0.16)
cx0 = Inches(0.6)
for i, (tag, name, period, body) in enumerate(ds):
    x = cx0 + i * (col_w + col_gap)
    # tag in mono
    text(s, x, Inches(4.85), col_w, Inches(0.3),
         tag, size=10, color=ACCENT, font=MONO, tracking=200)
    # name big
    text(s, x, Inches(5.15), col_w, Inches(0.5),
         name, size=22, bold=True, color=TEXT)
    # period
    text(s, x, Inches(5.7), col_w, Inches(0.3),
         period, size=11, italic=True, color=DIM)
    # accent rule
    rect(s, x, Inches(6.05), Inches(0.4), Emu(15000), fill=ACCENT)
    # facts
    lines(s, x, Inches(6.2), col_w, Inches(0.85),
          body[:3], size=10.5, color=DIM, line_space=Pt(2), font=MONO)
    # description italic
    text(s, x, Inches(6.85), col_w, Inches(0.3),
         body[3], size=10, italic=True, color=MUTED)
    # vertical sep
    if i < 2:
        vline(s, x + col_w + col_gap / 2, Inches(4.85), Inches(2.05), color=BORDER, width=Pt(0.5))


# =============================================================
# SLIDE 4 — Architecture
# =============================================================
s = prs.slides.add_slide(BLANK)
page_bg(s)
slide_meta(s, 4, "Architecture")

slide_title(s, "Three agents. One toolset.",
            "BaseAgent.analyze(question, dataset)  →  AnalysisResult")

# Horizontal flow:  INPUT → [3 agents stacked] → SHARED TOOLS

# Input column (left)
ix = Inches(0.6); iy = Inches(3.4); iw = Inches(2.0); ih = Inches(1.2)
rect(s, ix, iy, iw, ih, fill=SURFACE, line=BORDER, line_w=Pt(0.5))
rect(s, ix, iy, Emu(28000), ih, fill=ACCENT)
text(s, ix + Inches(0.25), iy + Inches(0.15), iw - Inches(0.3), Inches(0.3),
     "INPUT", size=9, color=ACCENT, font=MONO, tracking=200)
text(s, ix + Inches(0.25), iy + Inches(0.45), iw - Inches(0.3), Inches(0.4),
     "Question", size=15, bold=True, color=TEXT)
text(s, ix + Inches(0.25), iy + Inches(0.78), iw - Inches(0.3), Inches(0.3),
     "+ dataset", size=12, color=DIM)

# Three agents (middle, stacked)
agents = [
    ("Single",     "ReAct loop",                       "One LLM iterates with tools."),
    ("Multi",      "Planner › Analyst › Reviewer",     "Plan, execute, verify, retry."),
    ("RAG",        "Retrieve + Analyze",               "Codebook context grounds prompt."),
]
ax = Inches(3.5); aw = Inches(4.0); ah = Inches(1.05); a_gap = Inches(0.15)
ay0 = Inches(2.95)
for i, (name, sub, body) in enumerate(agents):
    ay = ay0 + i * (ah + a_gap)
    rect(s, ax, ay, aw, ah, fill=SURFACE, line=BORDER, line_w=Pt(0.5))
    rect(s, ax, ay, Emu(28000), ah, fill=ACCENT)
    text(s, ax + Inches(0.3), ay + Inches(0.12), aw - Inches(0.4), Inches(0.4),
         name, size=18, bold=True, color=TEXT)
    text(s, ax + Inches(0.3), ay + Inches(0.5), aw - Inches(0.4), Inches(0.3),
         sub, size=10.5, color=ACCENT, font=MONO)
    text(s, ax + Inches(0.3), ay + Inches(0.75), aw - Inches(0.4), Inches(0.3),
         body, size=11, color=DIM)
    # connector from input to agent
    line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                  ix + iw, iy + ih / 2,
                                  ax, ay + ah / 2)
    line.line.color.rgb = BORDER
    line.line.width = Pt(0.75)

# Shared tools column (right)
tx = Inches(8.1); ty = Inches(2.95); tw = Inches(4.65); th = Inches(3.5)
rect(s, tx, ty, tw, th, fill=SURFACE, line=BORDER, line_w=Pt(0.5))
text(s, tx + Inches(0.3), ty + Inches(0.2), tw - Inches(0.4), Inches(0.3),
     "SHARED TOOLS", size=10, color=ACCENT, font=MONO, tracking=200)
text(s, tx + Inches(0.3), ty + Inches(0.5), tw - Inches(0.4), Inches(0.4),
     "src/agents/tools.py", size=14, italic=True, color=TEXT, font=MONO)
hline(s, tx + Inches(0.3), ty + Inches(1.0), Inches(0.5), color=ACCENT, width=Pt(1.5))

tools = [
    ("load_dataset",        "smart loader · column subsetting"),
    ("get_dataset_schema",  "variable types & summary"),
    ("get_variable_info",   "codebook lookup"),
    ("run_analysis_code",   "sandboxed exec · timeout"),
]
for i, (name, desc) in enumerate(tools):
    yy = ty + Inches(1.25) + i * Inches(0.5)
    text(s, tx + Inches(0.3), yy, Inches(2.0), Inches(0.3),
         name, size=11, bold=True, color=ACCENT, font=MONO)
    text(s, tx + Inches(2.2), yy, Inches(2.4), Inches(0.3),
         desc, size=10, color=DIM)

# Connector from agents area to tools
line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                              ax + aw, Inches(4.2),
                              tx, Inches(4.2))
line.line.color.rgb = BORDER
line.line.width = Pt(0.75)


# =============================================================
# SLIDE 5 — Methodology
# =============================================================
s = prs.slides.add_slide(BLANK)
page_bg(s)
slide_meta(s, 5, "Methodology")

slide_title(s, "How we measure good analysis.",
            "30 questions  ×  3 agents  ×  3 datasets  =  90 controlled runs.")

# Big stat row
stats = [
    ("90", "TOTAL RUNS",         "3 × 30"),
    ("30", "QUESTIONS",          "10 per dataset"),
    ("4", "QUESTION TYPES",      "numeric · categorical · directional · descriptive"),
    ("T = 0", "TEMPERATURE",     "Claude Sonnet 4.5"),
]
sx0 = Inches(0.6); sw = Inches(2.95); sgap = Inches(0.13)
for i, (val, lbl, sub) in enumerate(stats):
    x = sx0 + i * (sw + sgap)
    text(s, x, Inches(2.85), sw, Inches(0.95),
         val, size=44, bold=True, color=ACCENT, font=MONO)
    text(s, x, Inches(3.8), sw, Inches(0.3),
         lbl, size=10, color=DIM, font=MONO, tracking=300)
    text(s, x, Inches(4.1), sw, Inches(0.3),
         sub, size=10, color=MUTED, italic=True)

# Divider
hline(s, Inches(0.6), Inches(4.85), Inches(12.13), color=BORDER, width=Pt(0.5))

# Two columns: design + scoring
text(s, Inches(0.6), Inches(5.05), Inches(6), Inches(0.3),
     "QUESTION DESIGN", size=10, color=ACCENT, font=MONO, tracking=200)
lines(s, Inches(0.6), Inches(5.4), Inches(6), Inches(1.6),
      ["Mix of weighted means, modal categories, correlations, and trends.",
       "Designed to require codebook lookup + survey-weight application.",
       "Identical question set across all three agents."],
      size=12, color=DIM, line_space=Pt(7))

text(s, Inches(7.0), Inches(5.05), Inches(5.7), Inches(0.3),
     "SCORING METRICS", size=10, color=ACCENT, font=MONO, tracking=200)
lines(s, Inches(7.0), Inches(5.4), Inches(5.7), Inches(1.6),
      ["Accuracy  ›  match against ground truth.",
       "Completeness  ›  addresses every part of the question.",
       "Weight Usage  ›  was the survey weight applied?",
       "Latency + retries tracked per run."],
      size=12, color=DIM, line_space=Pt(7))

# Vertical separator
vline(s, Inches(6.78), Inches(5.0), Inches(1.95), color=BORDER, width=Pt(0.5))


# =============================================================
# SLIDE 6 — Results
# =============================================================
s = prs.slides.add_slide(BLANK)
page_bg(s)
slide_meta(s, 6, "Results")

slide_title(s, "Architecture is the lever.",
            "Same LLM, same tools  ›  +37 points of accuracy from architecture alone.")

# Three giant headline numbers, side by side
headlines = [
    ("42%", "SINGLE",       "ReAct baseline",        DIM),
    ("62%", "RAG",          "+ codebook retrieval",  TEXT),
    ("79%", "MULTI-AGENT",  "plan · execute · verify", ACCENT),
]
hx0 = Inches(0.6); hw = Inches(3.95); hgap = Inches(0.16)
for i, (val, lbl, sub, color) in enumerate(headlines):
    x = hx0 + i * (hw + hgap)
    text(s, x, Inches(2.75), hw, Inches(1.4),
         val, size=88, bold=True, color=color, font=MONO)
    text(s, x, Inches(4.05), hw, Inches(0.4),
         lbl, size=12, bold=True, color=color, font=MONO, tracking=250)
    text(s, x, Inches(4.4), hw, Inches(0.3),
         sub, size=11, color=DIM, italic=True)

# Divider
hline(s, Inches(0.6), Inches(5.0), Inches(12.13), color=BORDER, width=Pt(0.5))

# Per-dataset breakdown — three small-multiples panels
text(s, Inches(0.6), Inches(5.2), Inches(8), Inches(0.3),
     "ACCURACY BY DATASET",
     size=10, color=ACCENT, font=MONO, tracking=200)

# Each panel: dataset header, then 3 horizontal bars with breathing room
panels = [
    ("GSS",            [("Single", 0.39, MUTED), ("RAG", 0.59, DIM), ("Multi", 0.60, ACCENT)]),
    ("WVS",            [("Single", 0.20, MUTED), ("RAG", 0.40, DIM), ("Multi", 0.79, ACCENT)]),
    ("ARAB BAROMETER", [("Single", 0.66, MUTED), ("RAG", 0.87, DIM), ("Multi", 0.98, ACCENT)]),
]
panel_w = Inches(3.95)
panel_gap = Inches(0.16)
panel_x0 = Inches(0.6)
panel_y = Inches(5.65)

for pi, (ds_name, rows) in enumerate(panels):
    px = panel_x0 + pi * (panel_w + panel_gap)
    # dataset header
    text(s, px, panel_y, panel_w, Inches(0.3),
         ds_name, size=11, bold=True, color=TEXT, font=MONO, tracking=200)
    # bars
    label_w = Inches(0.95)
    val_w = Inches(0.55)
    bar_area_w = panel_w - label_w - val_w
    bar_h = Inches(0.14)
    row_gap = Inches(0.36)
    bars_y0 = panel_y + Inches(0.4)
    for ri, (rname, rval, rcolor) in enumerate(rows):
        ry = bars_y0 + ri * row_gap
        # label
        text(s, px, ry - Inches(0.04), label_w, Inches(0.3),
             rname, size=10.5, color=TEXT,
             anchor=MSO_ANCHOR.MIDDLE)
        # track
        rect(s, px + label_w, ry + Inches(0.05), bar_area_w, bar_h, fill=SURFACE)
        # filled portion
        bw = int(bar_area_w * rval)
        if bw > 0:
            rect(s, px + label_w, ry + Inches(0.05), bw, bar_h, fill=rcolor)
        # value
        text(s, px + label_w + bar_area_w + Inches(0.08), ry - Inches(0.04),
             val_w, Inches(0.3),
             f"{int(rval*100)}%",
             size=10.5, bold=True, color=rcolor, font=MONO,
             anchor=MSO_ANCHOR.MIDDLE)
    # vertical separator
    if pi < len(panels) - 1:
        vline(s, px + panel_w + panel_gap / 2, panel_y, Inches(1.5),
              color=BORDER, width=Pt(0.5))


# =============================================================
# SLIDE 7 — Trade-offs
# =============================================================
s = prs.slides.add_slide(BLANK)
page_bg(s)
slide_meta(s, 7, "Trade-offs")

slide_title(s, "Quality has a cost.",
            "Multi-Agent wins accuracy at  ~3×  the latency. Weight discipline scales with structure.")

# Two custom bar charts side by side
text(s, Inches(0.6), Inches(2.85), Inches(6), Inches(0.3),
     "AVG LATENCY  ·  SECONDS / QUESTION",
     size=10, color=ACCENT, font=MONO, tracking=200)
hbar_chart(s, Inches(0.6), Inches(3.2), Inches(6.0), Inches(2.0),
           labels=["Single", "RAG", "Multi-Agent"],
           values=[52.7, 54.0, 166.4],
           max_value=200,
           accent=ACCENT,
           value_suffix="s",
           value_fmt="{:.1f}")

text(s, Inches(7.0), Inches(2.85), Inches(6), Inches(0.3),
     "WEIGHT USAGE  ·  PERCENT OF QUESTIONS",
     size=10, color=ACCENT, font=MONO, tracking=200)
hbar_chart(s, Inches(7.0), Inches(3.2), Inches(5.7), Inches(2.0),
           labels=["Single", "RAG", "Multi-Agent"],
           values=[60, 70, 100],
           max_value=110,
           accent=ACCENT,
           value_suffix="%",
           value_fmt="{:.0f}")

# Divider
hline(s, Inches(0.6), Inches(5.55), Inches(12.13), color=BORDER, width=Pt(0.5))

# Three short summary columns
sums = [
    ("SINGLE", "Fast — but skips methodology and misreads coded variables."),
    ("RAG", "Codebook grounding boosts accuracy at near-zero latency cost."),
    ("MULTI-AGENT", "Best accuracy and methodology — ~3× slower; 0.83 retries avg."),
]
for i, (lbl, body) in enumerate(sums):
    x = Inches(0.6 + i * 4.1)
    text(s, x, Inches(5.85), Inches(3.9), Inches(0.3),
         lbl, size=10, color=ACCENT, font=MONO, tracking=200)
    text(s, x, Inches(6.15), Inches(3.9), Inches(0.85),
         body, size=12, color=DIM)
    if i < 2:
        vline(s, x + Inches(3.95), Inches(5.85), Inches(1.0), color=BORDER, width=Pt(0.5))


# =============================================================
# SLIDE 8 — Failure Analysis
# =============================================================
s = prs.slides.add_slide(BLANK)
page_bg(s)
slide_meta(s, 8, "Failures")

slide_title(s, "Where the agents break.",
            "Errors cluster around coded variables and long-horizon planning.")

# Two columns — failure modes vs what RAG/Multi fixes
text(s, Inches(0.6), Inches(2.85), Inches(6), Inches(0.3),
     "✕   FAILURE MODES",
     size=11, color=WARM, font=MONO, tracking=200, bold=True)
fail_items = [
    "Wrong variable chosen — coded names (Q201A_1) without codebook lookup.",
    "Survey weights silently dropped (Single-Agent on WVS).",
    "Multi-step questions abandoned mid-plan (gss_03 – gss_05).",
    "Categorical labels confused with numeric codes (gss_02 'married' vs 1).",
]
lines(s, Inches(0.6), Inches(3.3), Inches(5.9), Inches(3.0),
      fail_items, size=12, color=DIM, line_space=Pt(10))

vline(s, Inches(6.7), Inches(2.85), Inches(3.5), color=BORDER, width=Pt(0.5))

text(s, Inches(7.0), Inches(2.85), Inches(6), Inches(0.3),
     "✓   WHAT THE ARCHITECTURE FIXES",
     size=11, color=ACCENT, font=MONO, tracking=200, bold=True)
fix_items = [
    "RAG retrieval rescues categorical-label questions   ›   +19 pts vs Single.",
    "Multi-Agent Reviewer step catches missing weights   ›   60% → 100% usage.",
    "Retry loop (≤ 2) recovers wrong-first-attempt runs  ›   0.83 retries avg.",
    "Long-horizon trend questions remain the open challenge for all three.",
]
lines(s, Inches(7.0), Inches(3.3), Inches(5.7), Inches(3.0),
      fix_items, size=12, color=DIM, line_space=Pt(10))

# Bottom callout — single line takeaway
hline(s, Inches(0.6), Inches(6.3), Inches(12.13), color=BORDER, width=Pt(0.5))
text(s, Inches(0.6), Inches(6.5), Inches(11.5), Inches(0.5),
     "Verification beats raw model capability — the Reviewer step alone moves weight-usage from 60% to 100%.",
     size=14, italic=True, color=ACCENT)


# =============================================================
# SLIDE 9 — Demo & Stack
# =============================================================
s = prs.slides.add_slide(BLANK)
page_bg(s)
slide_meta(s, 9, "Demo & Stack")

slide_title(s, "End-to-end live demo.",
            "Natural-language question  →  verified, weighted analysis.")

# Pipeline as 6 nodes, dot-and-line minimalist style
steps = [
    ("01", "ASK",     "natural-language\nquestion"),
    ("02", "CONTEXT", "schema +\ncodebook (RAG)"),
    ("03", "PLAN",    "decompose\nthe analysis"),
    ("04", "RUN",     "sandboxed\nPython exec"),
    ("05", "REVIEW",  "verify weights\n& methodology"),
    ("06", "ANSWER",  "code · output\ninterpretation"),
]
n = len(steps)
total_w = Inches(12.13)
node_d = Inches(0.32)
y_line = Inches(3.4)
x0 = Inches(0.6)
spacing = (total_w - node_d) / (n - 1)

# horizontal connecting line
hline(s, x0 + node_d / 2, y_line + node_d / 2,
      total_w - node_d, color=BORDER, width=Pt(0.75))

for i, (num, head, body) in enumerate(steps):
    cx = x0 + spacing * i
    cy = y_line
    # node — outer ring + inner dot
    ring = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, node_d, node_d)
    fill_solid(ring, INK)
    ring.line.color.rgb = ACCENT
    ring.line.width = Pt(1.5)
    inner = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               cx + Inches(0.1), cy + Inches(0.1),
                               Inches(0.12), Inches(0.12))
    fill_solid(inner, ACCENT); inner.line.fill.background()
    # number above
    text(s, cx - Inches(0.6), cy - Inches(0.5), Inches(1.5), Inches(0.3),
         num, size=10, color=MUTED, font=MONO, tracking=200, align=PP_ALIGN.CENTER)
    # head below
    text(s, cx - Inches(0.9), cy + Inches(0.55), Inches(2.1), Inches(0.4),
         head, size=13, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    # body
    text(s, cx - Inches(0.9), cy + Inches(0.95), Inches(2.1), Inches(0.7),
         body, size=10, color=DIM, align=PP_ALIGN.CENTER)

# Divider
hline(s, Inches(0.6), Inches(5.45), Inches(12.13), color=BORDER, width=Pt(0.5))

# Stack — four clean columns, hairline-separated (matches deck pattern)
text(s, Inches(0.6), Inches(5.65), Inches(8), Inches(0.3),
     "STACK", size=10, color=ACCENT, font=MONO, tracking=300)

stack_cols = [
    ("AI / ORCHESTRATION", ["Claude Sonnet 4.5", "LangGraph", "LangChain"]),
    ("RETRIEVAL",          ["ChromaDB", "MiniLM-L6-v2", "pdfplumber"]),
    ("DATA",               ["pandas", "numpy", "pyreadstat"]),
    ("INTERFACE",          ["FastAPI", "React", "TypeScript"]),
]
col_w = Inches(2.95)
col_gap = Inches(0.13)
col_x0 = Inches(0.6)
col_y = Inches(6.05)

for ci, (gname, items) in enumerate(stack_cols):
    cx = col_x0 + ci * (col_w + col_gap)
    # group label
    text(s, cx, col_y, col_w, Inches(0.3),
         gname, size=10, color=MUTED, font=MONO, tracking=200)
    # accent rule
    rect(s, cx, col_y + Inches(0.32), Inches(0.3), Emu(15000), fill=ACCENT)
    # items
    lines(s, cx, col_y + Inches(0.45), col_w, Inches(0.9),
          items, size=11, color=TEXT, font=MONO, line_space=Pt(3))
    # vertical separator
    if ci < len(stack_cols) - 1:
        vline(s, cx + col_w + col_gap / 2, col_y, Inches(1.0),
              color=BORDER, width=Pt(0.5))


# =============================================================
# SLIDE 10 — Closing
# =============================================================
s = prs.slides.add_slide(BLANK)
page_bg(s)

# Watermark
text(s, Inches(9.5), Inches(0.6), Inches(3.5), Inches(2.0),
     "10", size=130, bold=True, color=SURFACE, font=MONO, align=PP_ALIGN.RIGHT)

# Top kicker
text(s, Inches(0.6), Inches(0.7), Inches(8), Inches(0.3),
     "10  ·  CONCLUSION",
     size=10, color=ACCENT, font=MONO, tracking=300)

# Big closing statement
text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(1.0),
     "Architecture matters more",
     size=58, bold=True, color=TEXT)
text(s, Inches(0.6), Inches(3.25), Inches(12), Inches(1.0),
     "than model scale.",
     size=58, bold=True, color=ACCENT)

text(s, Inches(0.6), Inches(4.4), Inches(12), Inches(0.5),
     "Same LLM. Same tools. +37 points of accuracy from architecture alone.",
     size=18, italic=True, color=DIM)

# Two columns: takeaways + future
hline(s, Inches(0.6), Inches(5.4), Inches(12.13), color=BORDER, width=Pt(0.5))

text(s, Inches(0.6), Inches(5.55), Inches(6), Inches(0.3),
     "TAKEAWAYS", size=10, color=ACCENT, font=MONO, tracking=300)
lines(s, Inches(0.6), Inches(5.9), Inches(5.9), Inches(1.4),
      ["Architecture > model: +37 pts on identical LLM.",
       "Verification (Reviewer step) is the largest single accuracy lever.",
       "Codebook retrieval (RAG) gives +20 pts at near-zero latency cost."],
      size=11.5, color=DIM, line_space=Pt(6))

vline(s, Inches(6.78), Inches(5.5), Inches(1.5), color=BORDER, width=Pt(0.5))

text(s, Inches(7.0), Inches(5.55), Inches(6), Inches(0.3),
     "NEXT", size=10, color=ACCENT, font=MONO, tracking=300)
lines(s, Inches(7.0), Inches(5.9), Inches(5.7), Inches(1.4),
      ["Hybrid Multi-Agent + RAG.",
       "Cross-model study  ·  Claude vs GPT vs open-weights.",
       "Long-horizon planner & cost-aware routing."],
      size=11.5, color=DIM, line_space=Pt(6))

# Thank you + page
hline(s, Inches(0.6), Inches(7.0), Inches(12.13), color=BORDER, width=Pt(0.5))
text(s, Inches(0.6), Inches(7.1), Inches(8), Inches(0.3),
     "THANK YOU  ·  QUESTIONS?",
     size=10, color=ACCENT, font=MONO, tracking=400, bold=True)
text(s, Inches(11.0), Inches(7.1), Inches(1.7), Inches(0.3),
     "10 / 10", size=9, color=MUTED, font=MONO, align=PP_ALIGN.RIGHT)


# ---------- Save ----------
out = "/Users/sayed/Development/LLM Data Scientist/LLM_Data_Scientist_Expo.pptx"
prs.save(out)
print(f"Saved: {out}")
